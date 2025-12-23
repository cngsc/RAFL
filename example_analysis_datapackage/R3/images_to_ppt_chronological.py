''' README FIRST!! 
(2023-07-17)
 
updated code now sorts it according to chronological order (so 1-1 comes before 1-2 etc) and maintains aspect ratio of 
images while ensuring that image fits screen.
This code should be copied into your experiment folder (eg CN-BMC-304). 
In this folder, all your processed images should be saved in a folder titled "processed_images". 
Else, you can change the name from processed_images to whatever folder name your processed images are in (line27)
run code by opening terminal in folder (i suggest downloading vscode for ease), then typing: 

if mac : python3 images_to_ppt.py
if windows, change line31 "/" to "\", save it then run it: python images_to_ppt.py

Code will create a ppt file, one image per slide, with filename as slide title so you can quickly scroll through your images

'''

import os
import time
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# Get the date at which you run this
timestr = time.strftime("%Y%m%d-%H%M%S")
date = timestr.split("-")[0]  # Extract the date

# Get the path to the folder that contains the input images
input_folder_path = os.path.join(os.getcwd(), "processed_images")

# this command gets the folder that you have copied this code file into
path1 = os.getcwd()
# Specify the experiment name
expt = path1.split("/")[-1]

# Create a new PowerPoint presentation
presentation = Presentation()

# Create the initial title slide
slide_layout = presentation.slide_layouts[0]  # Presentation title slide
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = expt
slide.placeholders[1].text = date

# Create a dictionary to store the slide titles and filenames
slide_titles = {}

# Loop through each file in the folder
for filename in os.listdir(input_folder_path):
    if filename.endswith((".jpg", ".jpeg", ".png")):
        # Extract the slide title from the filename
        slide_title = os.path.splitext(filename)[0]

        # Store the slide title and filename in the dictionary
        slide_titles[slide_title] = filename

# Sort the slide titles in chronological order based on the first value before the hyphen
sorted_titles = sorted(slide_titles.keys(), key=lambda x: (int(x.split("-")[0]) if x.split("-")[0].isdigit() else 0, x))

# Create slides in chronological order
for slide_title in sorted_titles:
    filename = slide_titles[slide_title]

    # Create a new slide
    slide_layout = presentation.slide_layouts[5]  # Title and content
    slide = presentation.slides.add_slide(slide_layout)

    # Add the image to the slide
    left = Inches(0.5)
    top = Inches(1)
    target_width = Inches(9)
    target_height = Inches(6)

    # Load the image
    image_path = os.path.join(input_folder_path, filename)
    image = Image.open(image_path)

    # Calculate the image dimensions while maintaining the aspect ratio
    image_width, image_height = image.size
    aspect_ratio = min(target_width / image_width, target_height / image_height)
    new_width = int(image_width * aspect_ratio * 0.9)
    new_height = int(image_height * aspect_ratio * 0.9)

    # Calculate the position to center the image
    left += (target_width - new_width) / 2
    top += (target_height - new_height) / 2

    # Resize and add the image to the slide
    slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)

    # Set the slide title
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = slide_title

# Save the presentation
presentation.save("%s_ppt_%s.pptx" % (expt, date))
print("Presentation created successfully!")